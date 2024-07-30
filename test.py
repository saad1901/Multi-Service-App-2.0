import streamlit as st
import requests
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import anthropic

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_word(docx_file):
    document = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in document.paragraphs])
    return text

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_word(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_ppt(file)
    else:
        return "Unsupported file format."

# Streamlit app interface
st.title("Document Q&A with Anthropic")
st.write("Upload a document (PDF, Word, PowerPoint) and ask a question or give a command to summarize the content.")

anthropic_api_key = st.text_input("Anthropic API Key", type="password")
uploaded_file = st.file_uploader("Upload a document", type=["pdf", "docx", "pptx"])
prompt = st.text_input("Enter your prompt", placeholder="Summarize the document")

if uploaded_file and prompt and anthropic_api_key:
    with st.spinner('Processing...'):
        # Extract text from the uploaded document
        article = extract_text(uploaded_file)
        
        # Prepare the prompt for the API
        prompt_text = f"{anthropic.HUMAN_PROMPT} Here's a document:\n\n{article}\n\n\n\n{prompt}{anthropic.AI_PROMPT}"
        
        try:
            # Set the headers and data for the request
            headers = {
                "x-api-key": anthropic_api_key,
                "Content-Type": "application/json",
            }
            data = {
                "prompt": prompt_text,
                "model": "claude-v1",  # or "claude-2" if available
                "max_tokens": 100,  # Adjust as needed
                "stop_sequences": [anthropic.HUMAN_PROMPT],
            }
            
            # Make the POST request to the API
            response = requests.post("https://api.anthropic.com/v1/messages", headers=headers, json=data)
            
            if response.status_code == 200:
                result = response.json()
                # Display the response
                st.write("### Response")
                st.write(result.get('completion', 'No completion found'))  # Adjust based on response structure
            else:
                st.error(f"Error: {response.status_code} - {response.text}")
        except Exception as e:
            st.error(f"An error occurred: {e}")

elif not anthropic_api_key:
    st.info("Please enter your Anthropic API key.")

elif not uploaded_file:
    st.info("Please upload a document.")

elif not prompt:
    st.info("Please enter a prompt.")
